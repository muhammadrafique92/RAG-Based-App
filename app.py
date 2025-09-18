import os
import faiss
import streamlit as st
import numpy as np
from PyPDF2 import PdfReader
from docx import Document
import openpyxl
from pptx import Presentation
from sentence_transformers import SentenceTransformer
from groq import Groq

# ---------------- CONFIG ----------------
GROQ_API_KEY = os.getenv("GROQ_API_KEY")  # set this in Hugging Face Secrets
EMBED_MODEL = "all-MiniLM-L6-v2"          # free sentence transformer
VECTOR_DIM = 384                          # embedding size of all-MiniLM-L6-v2

# Init Groq client
client = Groq(api_key=GROQ_API_KEY)

# Load embedding model
embedder = SentenceTransformer(EMBED_MODEL)

# ---------------- HELPERS ----------------
def read_pdf(file):
    pdf = PdfReader(file)
    return " ".join([page.extract_text() for page in pdf.pages if page.extract_text()])

def read_docx(file):
    doc = Document(file)
    return " ".join([p.text for p in doc.paragraphs if p.text.strip()])

def read_xlsx(file):
    wb = openpyxl.load_workbook(file, data_only=True)
    text = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        for row in ws.iter_rows(values_only=True):
            row_text = " ".join([str(cell) for cell in row if cell])
            if row_text.strip():
                text.append(row_text)
    return " ".join(text)

def read_pptx(file):
    prs = Presentation(file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text.append(shape.text)
    return " ".join(text)

def chunk_text(text, chunk_size=500):
    words = text.split()
    return [" ".join(words[i:i+chunk_size]) for i in range(0, len(words), chunk_size)]

def embed_texts(texts):
    return embedder.encode(texts, convert_to_numpy=True)

def build_faiss_index(texts):
    vectors = embed_texts(texts)
    index = faiss.IndexFlatL2(VECTOR_DIM)
    index.add(vectors)
    return index, vectors

def retrieve(query, texts, vectors, index, k=3):
    q_vec = embed_texts([query])
    D, I = index.search(q_vec, k)
    return [texts[i] for i in I[0]]

def ask_groq(context, query):
    # hard limit context to ~2000 words (~3000 tokens)
    words = context.split()
    if len(words) > 2000:
        context = " ".join(words[:2000])

    messages = [
        {"role": "system", "content": "You are a helpful assistant. Use the context provided to answer."},
        {"role": "user", "content": f"Context:\n{context}\n\nQuestion: {query}"}
    ]
    completion = client.chat.completions.create(
        model="llama-3.3-70b-versatile",
        messages=messages,
        temperature=0.3,
        max_completion_tokens=512,
        top_p=1,
        stream=False
    )
    return completion.choices[0].message.content

# ---------------- STREAMLIT UI ----------------
st.set_page_config(page_title="RAG App", page_icon="📄", layout="centered")
st.title("📄 RAG App with Groq + Sentence Transformers")

uploaded_files = st.file_uploader(
    "Upload documents (PDF, DOCX, PPTX, XLSX)",
    type=["pdf", "docx", "pptx", "xlsx"],
    accept_multiple_files=True
)

if uploaded_files:
    texts = []
    for f in uploaded_files:
        if f.type == "application/pdf":
            content = read_pdf(f)
        elif f.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            content = read_docx(f)
        elif f.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            content = read_pptx(f)
        elif f.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            content = read_xlsx(f)
        else:
            content = ""

        # break long content into smaller chunks
        if content:
            texts.extend(chunk_text(content, chunk_size=500))

    if texts:
        st.success(f"✅ Loaded {len(texts)} chunks from documents")

        # Build FAISS index
        index, vectors = build_faiss_index(texts)

        query = st.text_input("🔍 Ask a question about your documents:")
        if query:
            with st.spinner("Retrieving and generating answer..."):
                retrieved = retrieve(query, texts, vectors, index, k=3)
                context = "\n".join(retrieved)
                answer = ask_groq(context, query)
            st.markdown("### 📝 Answer")
            st.write(answer)
    else:
        st.error("No readable text found in the uploaded documents.")
